from io import BytesIO
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


def export_xlsx(data) -> bytes:
    overall = data.get('overall', {})
    metrics = data.get('metrics', {})

    buf = BytesIO()
    with pd.ExcelWriter(buf, engine='openpyxl') as writer:
        pd.DataFrame([overall]).to_excel(writer, sheet_name='Overall', index=False)
        pd.DataFrame([metrics]).to_excel(writer, sheet_name='Metrics', index=False)
    buf.seek(0)
    return buf.getvalue()


def export_pptx(data) -> bytes:
    overall = data.get('overall', {})
    metrics = data.get('metrics', {})

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = 'Website Audit Report'
    slide.placeholders[1].text = f"Score: {overall.get('score','-')} | Grade: {overall.get('grade','-')} | Coverage: {overall.get('coverage','-')}%"

    # Metrics slide
    layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = 'Metrics'
    body = slide2.placeholders[1].text_frame
    body.clear()
    for k, v in metrics.items():
        p = body.add_paragraph()
        p.text = f"{k}: {v}"
        p.level = 0

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
