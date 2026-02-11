# -*- coding: utf-8 -*-

import io
import os
import re
import json
import time
import hashlib
import datetime as dt
from typing import Dict, Any, List, Tuple, Optional
from urllib.parse import urlparse, urljoin

import requests
from bs4 import BeautifulSoup

from reportlab.lib import colors
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
    PageBreak,
    Image
)
from reportlab.lib.units import inch
from reportlab.lib.pagesizes import A4
from reportlab.platypus.tableofcontents import TableOfContents
from reportlab.graphics.barcode import qr
from reportlab.graphics.shapes import Drawing
from reportlab.pdfgen import canvas

import matplotlib
matplotlib.use("Agg")  # Safe for headless servers/containers
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches


# =========================================================
# CONFIGURATION
# =========================================================

# NOTE: We preserve the original WEIGHTAGE keys to avoid breaking callers.
# Additional categories like "traffic" and "mobile" are supported for display
# but are not included in the weighted overall score unless you add them here.
WEIGHTAGE = {
    "performance": 0.30,
    "security": 0.25,
    "seo": 0.20,
    "accessibility": 0.15,
    "ux": 0.10,
}

PRIMARY_OK = colors.HexColor("#2ecc71")   # Green
PRIMARY_WARN = colors.HexColor("#f39c12") # Orange
PRIMARY_BAD = colors.HexColor("#e74c3c")  # Red
GRID = colors.HexColor("#DDE1E6")

DEFAULT_UA = (
    "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 "
    "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
)


# =========================================================
# WHITE LABEL BRANDING
# =========================================================

def get_branding(client_config: Dict[str, Any]):
    return {
        "company_name": client_config.get("company_name", "FF Tech"),
        "primary_color": client_config.get("primary_color", "#2c3e50"),
        "logo_path": client_config.get("logo_path", None)
    }


# =========================================================
# OPTIONAL LIGHTHOUSE (SAFE FALLBACK)
# =========================================================

def fetch_lighthouse_data(url: str, api_key: str = None) -> Dict[str, Any]:
    """
    Fetch minimal Lighthouse-like scores from PageSpeed Online.
    If not available or fails, return zeros. This function is optional and
    never blocks PDF generation.
    """
    try:
        endpoint = "https://www.googleapis.com/pagespeedonline/v5/runPagespeed"
        params = {"url": url, "strategy": "desktop"}
        if api_key:
            params["key"] = api_key
        resp = requests.get(endpoint, params=params, timeout=20)
        data = resp.json()
        cats = data.get("lighthouseResult", {}).get("categories", {})
        audits = data.get("lighthouseResult", {}).get("audits", {}) or {}
        lcp = audits.get("largest-contentful-paint", {}).get("numericValue")
        cls = audits.get("cumulative-layout-shift", {}).get("numericValue")
        inp = audits.get("interactive", {}).get("numericValue")
        return {
            "performance": cats.get("performance", {}).get("score", 0) * 100,
            "seo": cats.get("seo", {}).get("score", 0) * 100,
            "accessibility": cats.get("accessibility", {}).get("score", 0) * 100,
            "lcp_ms": float(lcp) if lcp is not None else None,
            "cls": float(cls) if cls is not None else None,
            "inp_ms": float(inp) if inp is not None else None,
        }
    except Exception:
        return {k: 0 for k in WEIGHTAGE}


# =========================================================
# BASIC APP-SIDE CHECKS (REAL-WORLD SAFE HEURISTICS)
# =========================================================

def run_basic_vulnerability_scan(url: str) -> List[str]:
    findings: List[str] = []
    try:
        if not url:
            return ["No URL provided for security scan"]
        r = requests.get(url, timeout=12, headers={"User-Agent": DEFAULT_UA})
        # Header checks:
        if "X-Frame-Options" not in r.headers:
            findings.append("Missing X-Frame-Options header")
        if "Content-Security-Policy" not in r.headers:
            findings.append("Missing Content-Security-Policy header")
        if "Strict-Transport-Security" not in r.headers and url.startswith("https"):
            findings.append("Missing HSTS header (Strict-Transport-Security)")

        # DOM checks:
        soup = BeautifulSoup(r.text, "html.parser")
        forms = soup.find_all("form")
        for f in forms:
            if not f.get("method"):
                findings.append("Form without method attribute detected")
        # Accessibility quick checks:
        imgs = soup.find_all("img")
        if imgs:
            missing_alt = sum(1 for im in imgs if not im.get("alt"))
            if missing_alt > 0:
                findings.append(f"{missing_alt} image(s) missing alt text)")
        # SEO quick checks:
        title_tag = soup.find("title")
        if not title_tag or not title_tag.get_text(strip=True):
            findings.append("Missing or empty <title> tag")
        meta_desc = soup.find("meta", attrs={"name": "description"})
        if not meta_desc or not meta_desc.get("content"):
            findings.append("Missing meta description")

        # robots.txt / sitemap hint
        try:
            robots = requests.get(_robots_url(url), timeout=6, headers={"User-Agent": DEFAULT_UA})
            if robots.status_code == 200:
                if "Sitemap:" not in robots.text:
                    findings.append("robots.txt found but no Sitemap directive present")
            else:
                findings.append("robots.txt not found or unreachable")
        except Exception:
            findings.append("robots.txt check failed")

    except Exception:
        findings.append("Scan failed or website unreachable")

    return findings


def _robots_url(url: str) -> str:
    try:
        p = urlparse(url)
        return f"{p.scheme}://{p.netloc}/robots.txt"
    except Exception:
        return url.rstrip("/") + "/robots.txt"


# =========================================================
# LIVE AUDIT (Fill sections from the actual site)
# =========================================================

def _fetch(url: str, timeout: int = 12) -> requests.Response:
    return requests.get(url, timeout=timeout, headers={"User-Agent": DEFAULT_UA}, allow_redirects=True)

def _safe_domain(url: str) -> str:
    try:
        return urlparse(url).netloc
    except Exception:
        return ""

def _is_same_domain(root: str, link: str) -> bool:
    try:
        rp, lp = urlparse(root), urlparse(link)
        return (rp.netloc.lower() == lp.netloc.lower()) and (lp.scheme in ("http", "https"))
    except Exception:
        return False

def _detect_cdn(headers: Dict[str, str]) -> str:
    hl = {k.lower(): v for k, v in headers.items()}
    if "cf-ray" in hl or "cf-cache-status" in hl:
        return "Cloudflare"
    if "x-amz-cf-id" in hl or "x-amz-cf-pop" in hl:
        return "AWS CloudFront"
    if "x-fastly-request-id" in hl:
        return "Fastly"
    if "x-akamai-transformed" in hl or "akamai-grn" in hl:
        return "Akamai"
    if "x-cache" in hl and "hit" in str(hl.get("x-cache", "")).lower():
        return "CDN (cache hit)"
    return "Unknown/No signal"

def _visible_text(soup: BeautifulSoup) -> str:
    for tag in soup(["script", "style", "noscript", "template"]):
        tag.extract()
    text = soup.get_text(separator=" ")
    text = re.sub(r"\s+", " ", text).strip()
    return text

def audit_live_site(url: str, link_sample: int = 12) -> Dict[str, Any]:
    result: Dict[str, Any] = {"ok": False}
    if not url:
        return result
    try:
        r = _fetch(url, timeout=18)
        ttfb_ms = r.elapsed.total_seconds() * 1000.0
        page_weight_mb = round(len(r.content or b"") / (1024 * 1024), 2)
        compressed = r.headers.get("content-encoding") in ("gzip", "br", "deflate")
        cdn = _detect_cdn(r.headers)
        content_type = r.headers.get("content-type", "")

        soup = BeautifulSoup(r.text, "html.parser")
        title = (soup.find("title").get_text(strip=True) if soup.find("title") else "")
        meta_desc_tag = soup.find("meta", attrs={"name": "description"})
        meta_desc_present = "Yes" if (meta_desc_tag and meta_desc_tag.get("content")) else "No"
        canonical = soup.find("link", rel=lambda x: x and "canonical" in x)
        canonical_url = canonical.get("href") if canonical else ""
        viewport = soup.find("meta", attrs={"name": "viewport"})
        has_viewport = True if viewport and ("width" in (viewport.get("content") or "").lower()) else False

        h1 = len(soup.find_all("h1"))
        h2 = len(soup.find_all("h2"))
        h3 = len(soup.find_all("h3"))

        imgs = soup.find_all("img")
        img_alt_coverage = round(100.0 * (sum(1 for im in imgs if (im.get("alt") not in (None, ""))) / max(1, len(imgs))), 2)

        aria_missing = 0
        for el in soup.find_all(True):
            if el.name in ("button", "a", "input") and not any(k.startswith("aria-") for k in el.attrs.keys()):
                aria_missing += 1

        schema_scripts = soup.find_all("script", attrs={"type": "application/ld+json"})
        has_schema = True if schema_scripts else False

        # Internal links sample for 404
        root = f"{urlparse(url).scheme}://{urlparse(url).netloc}"
        internal = []
        for a in soup.find_all("a", href=True):
            abs_url = urljoin(root, a.get("href"))
            if _is_same_domain(root, abs_url):
                internal.append(abs_url)
        sample = []
        seen = set()
        for l in internal:
            if l in seen:
                continue
            sample.append(l)
            seen.add(l)
            if len(sample) >= link_sample:
                break

        broken = 0
        for l in sample:
            try:
                rr = requests.head(l, timeout=6, headers={"User-Agent": DEFAULT_UA}, allow_redirects=True)
                st = rr.status_code
                if st >= 400:
                    rr2 = _fetch(l, timeout=8)
                    st = rr2.status_code
                if st >= 400:
                    broken += 1
            except Exception:
                broken += 1

        robots_present, sitemap_declared = False, False
        try:
            rob = requests.get(_robots_url(url), timeout=8, headers={"User-Agent": DEFAULT_UA})
            if rob.status_code == 200:
                robots_present = True
                sitemap_declared = ("sitemap:" in rob.text.lower())
        except Exception:
            pass

        # Security headers
        hl = {k.lower(): v for k, v in r.headers.items()}
        https_tls = url.lower().startswith("https")
        has_hsts = ("strict-transport-security" in hl)
        has_csp = ("content-security-policy" in hl)
        has_xfo = ("x-frame-options" in hl)
        has_xss = ("x-xss-protection" in hl)

        # Counts
        script_count = len(soup.find_all("script"))
        css_count = len(soup.find_all("link", attrs={"rel": "stylesheet"}))
        img_count = len(imgs)

        # Simple derived scores (0..100)
        perf = 100.0
        if ttfb_ms > 1500: perf -= 35
        elif ttfb_ms > 800: perf -= 20
        elif ttfb_ms > 400: perf -= 10
        if page_weight_mb > 6: perf -= 35
        elif page_weight_mb > 3: perf -= 20
        elif page_weight_mb > 1.5: perf -= 10
        perf -= max(0, (script_count + css_count + img_count) - 120) * 0.2
        if not compressed: perf -= 10
        if _detect_cdn(r.headers) == "Unknown/No signal": perf -= 5
        perf = max(0.0, min(100.0, perf))

        seo = 100.0
        if not title: seo -= 25
        if meta_desc_present != "Yes": seo -= 20
        if h1 != 1: seo -= 10
        if not robots_present: seo -= 10
        if not sitemap_declared: seo -= 5
        if broken > 0: seo -= min(25, broken * 3)
        if not canonical_url: seo -= 5
        if not has_schema: seo -= 5
        seo = max(0.0, min(100.0, seo))

        acc = 100.0
        if img_count > 0 and img_alt_coverage < 80: acc -= 25
        if aria_missing > 10: acc -= 15
        acc = max(0.0, min(100.0, acc))

        ux = 100.0
        if not has_viewport: ux -= 30
        if broken > 0: ux -= min(20, broken * 2)
        ux = max(0.0, min(100.0, ux))

        sec = 100.0
        if not https_tls: sec -= 40
        if https_tls and not has_hsts: sec -= 10
        if not has_csp: sec -= 20
        if not has_xfo: sec -= 10
        if not has_xss: sec -= 5
        sec = max(0.0, min(100.0, sec))

        result.update({
            "ok": True,
            "ttfb_ms": round(ttfb_ms, 2),
            "page_weight_mb": page_weight_mb,
            "compressed": "Yes" if compressed else "No",
            "cdn": cdn,
            "content_type": content_type,
            "title": title,
            "meta_description_present": meta_desc_present,
            "canonical_url": canonical_url,
            "viewport_present": "Yes" if has_viewport else "No",
            "h1_count": h1,
            "h2_count": h2,
            "h3_count": h3,
            "img_count": img_count,
            "img_alt_coverage_pct": img_alt_coverage,
            "aria_missing_count": aria_missing,
            "schema_org_present": "Yes" if has_schema else "No",
            "robots_present": "Yes" if robots_present else "No",
            "sitemap_declared": "Yes" if sitemap_declared else "No",
            "script_count": script_count,
            "css_count": css_count,
            "internal_links_sampled": len(sample),
            "broken_links_detected": broken,
            "security_headers": {
                "https_tls": "Yes" if https_tls else "No",
                "hsts": "Yes" if has_hsts else "No",
                "csp": "Yes" if has_csp else "No",
                "x_frame_options": "Yes" if has_xfo else "No",
                "x_xss_protection": "Yes" if has_xss else "No",
            },
            "scores": {
                "performance": perf,
                "seo": seo,
                "accessibility": acc,
                "ux": ux,
                "security": sec,
            }
        })

    except Exception:
        result["ok"] = False

    return result


# =========================================================
# SCORING
# =========================================================

def calculate_scores(audit_data: Dict[str, Any]):
    """
    Overall score is computed using WEIGHTAGE keys only (backward compatible).
    We also surface extra categories (traffic, mobile) in the Executive Summary
    if present in audit_data (but they don't affect overall unless added).
    """
    scores: Dict[str, float] = {}
    for k in WEIGHTAGE:
        val = float(audit_data.get(k, 0))
        scores[k] = max(0, min(val, 100))
    overall = sum(scores[k] * WEIGHTAGE[k] for k in WEIGHTAGE)
    return {
        "category_scores": scores,
        "overall_score": round(overall, 2),
    }


def score_to_status(score: float) -> Tuple[str, colors.Color]:
    if score >= 85:
        return "Good", PRIMARY_OK
    if score >= 65:
        return "Warning", PRIMARY_WARN
    return "Critical", PRIMARY_BAD


# =========================================================
# CHART HELPERS
# =========================================================

def _fig_to_buf() -> io.BytesIO:
    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format="png", dpi=150)
    plt.close()
    buf.seek(0)
    return buf


def line_chart(points: List[Tuple[str, float]], title: str, xlabel: str = "", ylabel: str = "") -> io.BytesIO:
    if not points:
        return _empty_chart("No data")
    labels = [p[0] for p in points]
    values = [float(p[1]) for p in points]
    fig, ax = plt.subplots(figsize=(6.4, 3.4))
    ax.plot(labels, values, marker="o", color="#0F62FE")
    ax.set_title(title)
    ax.set_xlabel(xlabel)
    ax.set_ylabel(ylabel)
    ax.grid(True, alpha=0.3)
    plt.xticks(rotation=25, ha="right")
    return _fig_to_buf()


def bar_chart(items: List[Tuple[str, float]], title: str, xlabel: str = "", ylabel: str = "") -> io.BytesIO:
    if not items:
        return _empty_chart("No data")
    labels = [i[0] for i in items]
    values = [float(i[1]) for i in items]
    fig, ax = plt.subplots(figsize=(6.4, 3.4))
    ax.bar(labels, values, color="#27AE60")
    ax.set_title(title)
    ax.set_xlabel(xlabel)
    ax.set_ylabel(ylabel)
    plt.xticks(rotation=25, ha="right")
    return _fig_to_buf()


def pie_chart(parts: List[Tuple[str, float]], title: str = "") -> io.BytesIO:
    if not parts:
        return _empty_chart("No data")
    labels = [p[0] for p in parts]
    sizes = [max(0.0, float(p[1])) for p in parts]
    if sum(sizes) <= 0:
        sizes = [1.0 for _ in sizes]
    fig, ax = plt.subplots(figsize=(5, 5))
    ax.pie(sizes, labels=labels, autopct="%1.1f%%", startangle=140)
    ax.axis("equal")
    ax.set_title(title)
    return _fig_to_buf()


def radar_chart(categories: List[str], values: List[float], title: str = "") -> io.BytesIO:
    if not categories or not values:
        return _empty_chart("No data")
    N = len(categories)
    vals = [float(v) for v in values[:N]]
    angles = np.linspace(0, 2 * np.pi, N, endpoint=False).tolist()
    vals += vals[:1]
    angles += angles[:1]
    fig = plt.figure(figsize=(5.6, 5.0))
    ax = fig.add_subplot(111, polar=True)
    ax.set_theta_offset(np.pi / 2)
    ax.set_theta_direction(-1)
    ax.set_thetagrids(np.degrees(angles[:-1]), categories)
    ax.set_ylim(0, 100)
    ax.plot(angles, vals, color="#2F80ED", linewidth=2)
    ax.fill(angles, vals, color="#2F80ED", alpha=0.2)
    ax.set_title(title, y=1.1)
    return _fig_to_buf()


def heatmap(matrix: List[List[float]], xlabels: List[str], ylabels: List[str], title: str = "") -> io.BytesIO:
    data = np.array(matrix) if matrix else np.zeros((1, 1))
    fig, ax = plt.subplots(figsize=(6.4, 3.8))
    c = ax.imshow(data, cmap("RdYlGn"), vmin=0, vmax=100, aspect="auto")
    ax.set_xticks(range(len(xlabels)))
    ax.set_yticks(range(len(ylabels)))
    ax.set_xticklabels(xlabels, rotation=35, ha="right")
    ax.set_yticklabels(ylabels)
    ax.set_title(title)
    fig.colorbar(c, ax=ax, fraction=0.046, pad=0.04)
    return _fig_to_buf()


def multi_line_chart(series: Dict[str, List[Tuple[str, float]]], title: str, xlabel: str = "", ylabel: str = "") -> io.BytesIO:
    if not series:
        return _empty_chart("No data")
    fig, ax = plt.subplots(figsize=(6.8, 3.8))
    for name, pts in series.items():
        labels = [p[0] for p in pts]
        values = [float(p[1]) for p in pts]
        ax.plot(labels, values, marker="o", linewidth=2, label=name)
    ax.set_title(title)
    ax.set_xlabel(xlabel)
    ax.set_ylabel(ylabel)
    ax.grid(True, alpha=0.3)
    plt.xticks(rotation=25, ha="right")
    ax.legend()
    return _fig_to_buf()


def _empty_chart(msg: str) -> io.BytesIO:
    fig, ax = plt.subplots(figsize=(5, 3))
    ax.axis("off")
    ax.text(0.5, 0.5, msg, ha="center", va="center")
    return _fig_to_buf()


# =========================================================
# DIGITAL SIGNATURE
# =========================================================

def generate_digital_signature(content: bytes) -> str:
    return hashlib.sha256(content).hexdigest()


# =========================================================
# POWERPOINT AUTO GENERATION (Optional)
# =========================================================

def generate_executive_ppt(audit_data: Dict[str, Any], file_path: str):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Executive Audit Summary"
    content = slide.placeholders[1]
    content.text = f"Overall Score: {audit_data.get('overall_score', 0)}"
    prs.save(file_path)


# =========================================================
# DOC WITH CLICKABLE TOC + HEADER/FOOTER
# =========================================================

class _DocWithTOC(SimpleDocTemplate):
    def __init__(self, *args, **kwargs):
        self._heading_styles = {"Heading1": 0, "Heading2": 1}
        super().__init__(*args, **kwargs)

    def afterFlowable(self, flowable):
        from reportlab.platypus import Paragraph
        if isinstance(flowable, Paragraph):
            style_name = getattr(flowable.style, "name", "")
            if style_name in self._heading_styles:
                level = self._heading_styles[style_name]
                text = flowable.getPlainText()
                key = f"h_{hash((text, self.page))}"
                self.canv.bookmarkPage(key)
                self.canv.addOutlineEntry(text, key, level=level, closed=False)
                self.notify("TOCEntry", (level, text, self.page))


def _draw_header_footer(c: canvas.Canvas, doc: SimpleDocTemplate, url: str, branding: Dict[str, Any]):
    c.saveState()
    # Header line
    c.setStrokeColor(colors.HexColor(branding.get("primary_color", "#2c3e50")))
    c.setLineWidth(0.6)
    c.line(doc.leftMargin, doc.height + doc.topMargin + 6, doc.width + doc.leftMargin, doc.height + doc.topMargin + 6)
    # Header text
    c.setFont("Helvetica-Bold", 9)
    c.setFillColor(colors.HexColor("#525252"))
    header_left = f"{branding.get('company_name', 'FF Tech')} — {url or ''}"
    c.drawString(doc.leftMargin, doc.height + doc.topMargin + 10, header_left[:140])
    # Footer
    c.setFont("Helvetica", 8)
    c.setFillColor(colors.HexColor("#6F6F6F"))
    page_num = c.getPageNumber()
    timestamp = dt.datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC")
    footer_left = f"Generated: {timestamp}"
    footer_right = f"Page {page_num}"
    c.drawString(doc.leftMargin, doc.bottomMargin - 20, footer_left)
    w = c.stringWidth(footer_right, "Helvetica", 8)
    c.drawString(doc.leftMargin + doc.width - w, doc.bottomMargin - 20, footer_right)
    c.restoreState()


# =========================================================
# MAIN GENERATOR (UNCHANGED SIGNATURE/RETURN)
# =========================================================

def generate_audit_pdf(audit_data: Dict[str, Any],
                       client_config: Dict[str, Any] = None,
                       history_scores: List[float] = None) -> bytes:

    branding = get_branding(client_config or {})
    url = audit_data.get("url", "")
    tool_version = audit_data.get("tool_version", "v1.0")

    # Run the live audit for real data filling
    live = audit_live_site(url)

    # If caller didn't supply category scores, derive from live audit
    derived_scores = live.get("scores", {}) if live.get("ok") else {}
    merged_for_overall = {
        k: float(audit_data.get(k, derived_scores.get(k, 0))) for k in WEIGHTAGE
    }
    score_data = {
        "category_scores": merged_for_overall,
        "overall_score": round(sum(merged_for_overall[k] * WEIGHTAGE[k] for k in WEIGHTAGE), 2)
    }

    buf = io.BytesIO()
    doc = _DocWithTOC(
        buf, pagesize=A4,
        leftMargin=36, rightMargin=36, topMargin=54, bottomMargin=54,
        title="FF Tech Web Audit",
        author=branding.get("company_name", "FF Tech")
    )
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="KPIHeader", parent=styles["Heading2"],
                              textColor=colors.HexColor(branding.get("primary_color", "#2c3e50"))))
    styles.add(ParagraphStyle(name="Small", parent=styles["Normal"], fontSize=9))
    styles.add(ParagraphStyle(name="Tiny", parent=styles["Normal"], fontSize=8, textColor=colors.HexColor("#6F6F6F")))
    elements: List[Any] = []

    # -------------------- 1) COVER PAGE --------------------
    logo_path = branding.get("logo_path")
    if logo_path and os.path.exists(logo_path):
        try:
            elements.append(Image(logo_path, width=1.6 * inch, height=1.6 * inch))
        except Exception:
            pass

    elements.append(Paragraph(branding["company_name"], styles["Title"]))
    elements.append(Spacer(1, 0.12 * inch))
    elements.append(Paragraph("FF Tech Web Audit Report", styles["Heading1"]))
    elements.append(Spacer(1, 0.06 * inch))
    elements.append(Paragraph(f"Website: {url}", styles["Normal"]))
    elements.append(Paragraph(f"Domain: {_safe_domain(url)}", styles["Small"]))
    elements.append(Paragraph(f"Report Time (UTC): {dt.datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')}", styles["Normal"]))
    elements.append(Paragraph(f"Audit/Tool Version: {tool_version}", styles["Normal"]))
    elements.append(Spacer(1, 0.12 * inch))

    dashboard_link = audit_data.get("dashboard_url") or url
    if dashboard_link:
        try:
            code = qr.QrCodeWidget(dashboard_link)
            bx = code.getBounds()
            w = bx[2] - bx[0]; h = bx[3] - bx[1]
            d = Drawing(1.4 * inch, 1.4 * inch, transform=[1.4 * inch / w, 0, 0, 1.4 * inch / h, 0, 0])
            d.add(code)
            elements.append(d)
            if "Caption" in styles:
                elements.append(Paragraph("Scan to view the online audit/dashboard", styles["Caption"]))
        except Exception:
            pass

    elements.append(PageBreak())

    # -------------------- 3) CLICKABLE TOC --------------------
    toc = TableOfContents()
    toc.levelStyles = [
        ParagraphStyle(fontName="Helvetica-Bold", name="TOCHeading1", fontSize=12, leftIndent=20, firstLineIndent=-10, spaceBefore=6),
        ParagraphStyle(fontName="Helvetica", name="TOCHeading2", fontSize=10, leftIndent=36, firstLineIndent=-10, spaceBefore=2),
    ]
    elements.append(Paragraph("Table of Contents", styles["Heading1"]))
    elements.append(Spacer(1, 0.08 * inch))
    elements.append(toc)
    elements.append(PageBreak())

    # -------------------- 2) EXECUTIVE SUMMARY --------------------
    elements.append(Paragraph("Executive Summary", styles["Heading1"]))
    elements.append(Spacer(1, 0.06 * inch))
    elements.append(Paragraph(f"Overall Website Health Score: <b>{score_data['overall_score']}</b>/100", styles["Normal"]))

    cat_rows = [["Category", "Score", "Status"]]
    for k in WEIGHTAGE:
        s = score_data["category_scores"].get(k, 0)
        status, _ = score_to_status(s)
        cat_rows.append([k.capitalize(), f"{s:.0f}", status])

    # Optional extra categories if caller provided
    if "traffic_score" in audit_data:
        ts = float(audit_data["traffic_score"]) if audit_data["traffic_score"] is not None else 0
        cat_rows.append(["Traffic & Engagement", f"{ts:.0f}", score_to_status(ts)[0]])
    if "mobile" in audit_data:
        ms = float(audit_data["mobile"]) if audit_data["mobile"] is not None else 0
        cat_rows.append(["Mobile Responsiveness", f"{ms:.0f}", score_to_status(ms)[0]])

    table = Table(cat_rows, colWidths=[3.1 * inch, 1.0 * inch, 1.6 * inch])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#F2F4F8")),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#121619")),
        ("GRID", (0, 0), (-1, -1), 0.25, GRID),
        ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("ALIGN", (1, 1), (-1, -1), "CENTER"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#FAFAFA")]),
    ]))
    elements.append(table)

    # Live highlights to ensure filled attributes
    if live.get("ok"):
        elements.append(Spacer(1, 0.10 * inch))
        hi_rows = [
            ["TTFB (ms)", str(live.get("ttfb_ms"))],
            ["Page Weight (MB)", str(live.get("page_weight_mb"))],
            ["Compression", live.get("compressed", "No")],
            ["CDN", live.get("cdn", "Unknown")],
            ["robots.txt", live.get("robots_present", "No")],
            ["Sitemap Declared", live.get("sitemap_declared", "No")],
            ["Broken Links (sample)", str(live.get("broken_links_detected", 0))],
        ]
        elements.append(_mini_table(["Highlight", "Value"], hi_rows))

    # Trend summary chart if history_scores provided (last 6–12 months)
    if history_scores:
        img = line_chart([(str(i + 1), float(v)) for i, v in enumerate(history_scores)], "Overall Score Trend", "Period", "Score")
        elements.append(Image(img, width=5.8 * inch, height=3.2 * inch))

    # AI recommendations summary
    elements.append(Spacer(1, 0.12 * inch))
    elements.append(Paragraph("AI-Generated Recommendations (Summary)", styles["KPIHeader"]))
    for rec in _collect_ai_recommendations(live, score_data["category_scores"]):
        elements.append(Paragraph(f"- {rec}", styles["Normal"]))

    elements.append(PageBreak())

    # -------------------- 4) TRAFFIC & GOOGLE SEARCH METRICS --------------------
    elements.append(Paragraph("Traffic & Google Search Metrics", styles["Heading1"]))
    traffic = audit_data.get("traffic") or {}
    gsc = audit_data.get("gsc") or {}
    if traffic:
        if isinstance(traffic.get("trend"), list) and traffic["trend"]:
            try:
                pts = [(str(p[0]), float(p[1])) for p in traffic["trend"]]
                img = line_chart(pts, "Traffic Trend (6–12 months)", "Period", "Visits")
                elements.append(Image(img, width=5.8 * inch, height=3.2 * inch))
            except Exception:
                pass
        if isinstance(traffic.get("sources"), dict) and traffic["sources"]:
            try:
                parts = [(k, float(v)) for k, v in traffic["sources"].items()]
                img = pie_chart(parts, "Traffic Sources")
                elements.append(Image(img, width=4.8 * inch, height=4.8 * inch))
            except Exception:
                pass
        ga_rows = []
        def _ga(label, key):
            v = traffic.get(key)
            ga_rows.append([label, str(v) if v is not None else "N/A"])
        for label, key in [
            ("Total Visitors", "total_visitors"),
            ("Organic Traffic", "organic"),
            ("Direct", "direct"),
            ("Referral", "referral"),
            ("Social", "social"),
            ("Paid", "paid"),
            ("Bounce Rate (%)", "bounce_rate"),
            ("Avg. Session Duration (s)", "avg_session_duration"),
            ("Pageviews per Session", "pages_per_session"),
        ]: _ga(label, key)
        elements.append(_mini_table(["Metric", "Value"], ga_rows))
    else:
        elements.append(Paragraph("Google Analytics/GA4 data not available (no credentials/data provided).", styles["Small"]))

    if gsc:
        gsc_rows = []
        for lbl, key in [("Impressions", "impressions"), ("Clicks", "clicks"), ("CTR (%)", "ctr")]:
            gsc_rows.append([lbl, str(gsc.get(key)) if gsc.get(key) is not None else "N/A"])
        elements.append(_mini_table(["GSC Metric", "Value"], gsc_rows))
        if isinstance(gsc.get("queries"), list) and gsc["queries"]:
            try:
                top_q = gsc["queries"][:10]
                items = [(str(q.get("query", ""))[:24], float(q.get("clicks", 0))) for q in top_q]
                img = bar_chart(items, "Top Queries by Clicks", "Query", "Clicks")
                elements.append(Image(img, width=5.8 * inch, height=3.2 * inch))
            except Exception:
                pass
    else:
        elements.append(Paragraph("Google Search Console data not available (no credentials/data provided).", styles["Small"]))

    elements.append(PageBreak())

    # -------------------- 5) SEO KPIs (~40)
    elements.append(Paragraph("SEO KPIs", styles["Heading1"]))
    seo_rows = _build_seo_kpis_table(url, audit_data, live)
    elements.append(_kpi_scorecard_table(seo_rows))
    elements.append(PageBreak())

    # -------------------- 6) PERFORMANCE KPIs (~20)
    elements.append(Paragraph("Performance KPIs", styles["Heading1"]))
    perf_rows = _build_performance_kpis_table(audit_data, live)
    elements.append(_kpi_scorecard_table(perf_rows))
    elements.append(PageBreak())

    # -------------------- 7) SECURITY KPIs (15–20)
    elements.append(Paragraph("Security KPIs", styles["Heading1"]))
    sec_rows = _build_security_kpis_table(url, audit_data, live)
    elements.append(_kpi_scorecard_table(sec_rows))
    findings = run_basic_vulnerability_scan(url)
    if findings:
        elements.append(Paragraph("Automated Quick Findings", styles["KPIHeader"]))
        for f in findings:
            elements.append(Paragraph(f"- {f}", styles["Normal"]))
    elements.append(PageBreak())

    # -------------------- 8) ACCESSIBILITY KPIs (15–20)
    elements.append(Paragraph("Accessibility KPIs", styles["Heading1"]))
    acc_rows = _build_accessibility_kpis_table(audit_data, live)
    elements.append(_kpi_scorecard_table(acc_rows))
    # Radar (derived)
    acc_radar = {
        "Alt Coverage": live.get("img_alt_coverage_pct", 0),
        "ARIA Coverage": max(0.0, 100.0 - min(100.0, live.get("aria_missing_count", 0) * 5.0)),
        "Structure (H1/H2/H3)": 100.0 if live.get("h1_count") == 1 else 70.0,
        "Viewport": 100.0 if live.get("viewport_present") == "Yes" else 40.0,
    }
    img = radar_chart(list(acc_radar.keys()), list(acc_radar.values()), "Accessibility Radar")
    elements.append(Image(img, width=5.2 * inch, height=4.8 * inch))
    elements.append(PageBreak())

    # -------------------- 9) UX / User Experience KPIs (15–20)
    elements.append(Paragraph("UX / User Experience KPIs", styles["Heading1"]))
    ux_rows = _build_ux_kpis_table(audit_data, live)
    elements.append(_kpi_scorecard_table(ux_rows))
    ux_radar = {
        "Mobile": 100.0 if live.get("viewport_present") == "Yes" else 40.0,
        "Broken Links": max(0.0, 100.0 - min(100.0, live.get("broken_links_detected", 0) * 8.0)),
        "Readability (proxy)": 70.0,
        "Interactivity (proxy)": 70.0,
    }
    img = radar_chart(list(ux_radar.keys()), list(ux_radar.values()), "UX Radar")
    elements.append(Image(img, width=5.2 * inch, height=4.8 * inch))
    elements.append(PageBreak())

    # -------------------- 10) Competitor Comparison
    elements.append(Paragraph("Competitor Comparison", styles["Heading1"]))
    competitors = audit_data.get("competitors") or []
    if competitors:
        series = {}
        for comp in competitors[:5]:
            name = comp.get("name") or comp.get("domain") or "Competitor"
            trend = comp.get("traffic_trend") or []
            series[name] = [(str(p[0]), float(p[1])) for p in trend][:12]
        if series:
            img = multi_line_chart(series, "Traffic vs Competitors", "Period", "Traffic")
            elements.append(Image(img, width=6.2 * inch, height=3.6 * inch))
        seo_comp = [(str((c.get("name") or c.get("domain") or "Comp"))[:22], float(c.get("seo", 0))) for c in competitors[:5]]
        if seo_comp:
            img = bar_chart(seo_comp, "SEO Performance (Score)", "Competitor", "Score")
            elements.append(Image(img, width=6.2 * inch, height=3.6 * inch))
    else:
        elements.append(Paragraph("Not available (no competitor data provided).", styles["Small"]))
    elements.append(PageBreak())

    # -------------------- 11) Historical Comparison / Trend Analysis
    elements.append(Paragraph("Historical Comparison / Trend Analysis", styles["Heading1"]))
    hist = audit_data.get("history") or {}
    for key, label, ylabel in [
        ("traffic", "Traffic Trend", "Traffic"),
        ("keyword_rank", "Keyword Ranking Trend", "Avg Rank"),
        ("page_speed", "Page Speed Trend (ms)", "ms"),
        ("sec_vulns", "Security Vulnerabilities Trend", "Count"),
        ("engagement", "Engagement Trend (Pages/Session)", "Pages/Session"),
    ]:
        if key in hist and isinstance(hist[key], list) and hist[key]:
            try:
                pts = [(str(p[0]), float(p[1])) for p in hist[key]][:12]
                img = line_chart(pts, label, "Period", ylabel)
                elements.append(Image(img, width=5.8 * inch, height=3.2 * inch))
            except Exception:
                pass
    elements.append(PageBreak())

    # -------------------- 13) KPI Scorecards (Consolidated)
    elements.append(Paragraph("KPI Scorecards (Consolidated)", styles["Heading1"]))
    consolidated = _consolidate_kpis(
        seo_rows,
        perf_rows,
        sec_rows,
        acc_rows,
        ux_rows,
    )
    elements.append(_kpi_scorecard_table(consolidated, show_weight=True, show_link=False))
    elements.append(PageBreak())

    # -------------------- 14) AI Recommendations (Detailed)
    elements.append(Paragraph("AI Recommendations (Detailed)", styles["Heading1"]))
    detailed_recs = _collect_ai_recommendations(live, score_data["category_scores"])
    for r in detailed_recs:
        elements.append(Paragraph(f"- {r}", styles["Normal"]))

    # Build document
    def _first(c, d): _draw_header_footer(c, d, url, branding)
    def _later(c, d): _draw_header_footer(c, d, url, branding)
    doc.build(elements, onFirstPage=_first, onLaterPages=_later)

    pdf_bytes = buf.getvalue()
    buf.close()

    # Signature (stdout/log)
    signature = generate_digital_signature(pdf_bytes)
    print("Digital Signature:", signature)

    # Optional PPT (kept for backward compatibility)
    try:
        generate_executive_ppt({"overall_score": score_data["overall_score"]}, "/tmp/executive_summary.pptx")
    except Exception:
        pass

    return pdf_bytes


# =========================================================
# HELPERS FOR KPI TABLES / RECOMMENDATIONS
# =========================================================

def _mini_table(headers: List[str], rows: List[List[Any]]) -> Table:
    data = [headers] + rows
    t = Table(data, colWidths=[2.6 * inch, 2.0 * inch])
    t.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#F2F4F8")),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("GRID", (0,0), (-1,-1), 0.25, GRID),
        ("FONTSIZE", (0,0), (-1,-1), 9),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
    ]))
    return t


def _kpi_scorecard_table(rows: List[Dict[str, Any]], show_weight: bool = False, show_link: bool = False) -> Table:
    headers = ["KPI Name", "Value", "Status"]
    col_widths = [3.1 * inch, 1.1 * inch, 1.0 * inch]
    if show_weight:
        headers.append("Weight")
        col_widths.append(0.9 * inch)
    if show_link:
        headers.append("Link")
        col_widths.append(1.3 * inch)

    data = [headers]
    for r in rows:
        name = str(r.get("name", ""))
        value = r.get("value", "")
        status = str(r.get("status", ""))
        weight = r.get("weight", "")
        link = r.get("link", "")
        row = [name, str(value), status]
        if show_weight:
            row.append(f"{weight}" if isinstance(weight, (int, float)) else str(weight or ""))
        if show_link:
            row.append(str(link)[:60] if link else "")
        data.append(row)

    t = Table(data, colWidths=col_widths, repeatRows=1)
    styles = [
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#F2F4F8")),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#121619")),
        ("GRID", (0, 0), (-1, -1), 0.25, GRID),
        ("FONTNAME", (0, 1), (-1, -1), "Helvetica"),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#FAFAFA")]),
    ]
    # Color-code status column
    for r_i in range(1, len(data)):
        try:
            status = (data[r_i][2] or "").strip().lower()
            if status.startswith("good"):
                styles.append(("TEXTCOLOR", (2, r_i), (2, r_i), PRIMARY_OK))
            elif status.startswith("warn)
                styles.append(("TEXTCOLOR", (2, r_i), (2, r_i), PRIMARY_WARN))
            else:
                styles.append(("TEXTCOLOR", (2, r_i), (2, r_i), PRIMARY_BAD))
        except Exception:
            pass
    t.setStyle(TableStyle(styles))
    return t


def _build_seo_kpis_table(url: str, audit_data: Dict[str, Any], live: Dict[str, Any]) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    custom = audit_data.get("seo_kpis") or []
    rows.extend(custom)

    # Live checks
    rows += [
        {"name": "Title Tag Present", "value": "Yes" if live.get("title") else "No", "status": "Good" if live.get("title") else "Critical"},
        {"name": "Meta Description Present", "value": live.get("meta_description_present", "No"), "status": "Good" if live.get("meta_description_present") == "Yes" else "Warning"},
        {"name": "Canonical Tag Present", "value": "Yes" if live.get("canonical_url") else "No", "status": "Good" if live.get("canonical_url") else "Warning"},
        {"name": "H1 Count", "value": live.get("h1_count", 0), "status": "Good" if live.get("h1_count", 0) == 1 else "Warning"},
        {"name": "H2 Count", "value": live.get("h2_count", 0), "status": "Good" if live.get("h2_count", 0) >= 1 else "Warning"},
        {"name": "H3 Count", "value": live.get("h3_count", 0), "status": "Good" if live.get("h3_count", 0) >= 1 else "Warning"},
        {"name": "robots.txt Present", "value": live.get("robots_present", "No"), "status": "Good" if live.get("robots_present") == "Yes" else "Warning"},
        {"name": "Sitemap Declared", "value": live.get("sitemap_declared", "No"), "status": "Good" if live.get("sitemap_declared") == "Yes" else "Warning"},
        {"name": "Broken Links (sample)", "value": live.get("broken_links_detected", 0), "status": "Good" if live.get("broken_links_detected", 0) == 0 else "Warning"},
        {"name": "Image Alt Coverage (%)", "value": live.get("img_alt_coverage_pct", 0), "status": "Good" if live.get("img_alt_coverage_pct", 0) >= 80 else "Warning"},
        {"name": "Structured Data (Schema.org)", "value": live.get("schema_org_present", "No"), "status": "Good" if live.get("schema_org_present") == "Yes" else "Warning"},
    ]

    return rows[:40]


def _build_performance_kpis_table(audit_data: Dict[str, Any], live: Dict[str, Any]) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    perf = audit_data.get("performance_kpis") or {}

    def add(name, value, good_when):
        status = "Good" if good_when(value) else "Warning"
        rows.append({"name": name, "value": value, "status": status})

    add("Time To First Byte (ms)", live.get("ttfb_ms", "N/A"), lambda v: isinstance(v, (int, float)) and v <= 800)
    add("Page Weight (MB)", live.get("page_weight_mb", "N/A"), lambda v: isinstance(v, (int, float)) and v <= 3)
    rows.append({"name": "Compression (gzip/brotli)", "value": live.get("compressed", "No"), "status": "Good" if live.get("compressed") == "Yes" else "Warning"})
    rows.append({"name": "CDN Detected", "value": live.get("cdn", "Unknown"), "status": "Good" if live.get("cdn") != "Unknown/No signal" else "Warning"})
    add("Scripts (count)", live.get("script_count", 0), lambda v: isinstance(v, (int, float)) and v <= 60)
    add("Stylesheets (count)", live.get("css_count", 0), lambda v: isinstance(v, (int, float)) and v <= 15)
    add("Images (count)", live.get("img_count", 0), lambda v: isinstance(v, (int, float)) and v <= 120)

    # Extra provided metrics
    for k, v in perf.items():
        label = k.replace("_", " ").title()
        if all(label != r["name"] for r in rows):
            rows.append({"name": label, "value": v, "status": "Good"})

    return rows[:20]


def _build_security_kpis_table(url: str, audit_data: Dict[str, Any], live: Dict[str, Any]) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    sh = live.get("security_headers", {}) or {}

    rows += [
        {"name": "HTTPS / TLS", "value": sh.get("https_tls", "No"), "status": "Good" if sh.get("https_tls") == "Yes" else "Critical"},
        {"name": "HSTS Header", "value": sh.get("hsts", "No"), "status": "Good" if sh.get("hsts") == "Yes" else "Warning"},
        {"name": "Content-Security-Policy", "value": sh.get("csp", "No"), "status": "Good" if sh.get("csp") == "Yes" else "Warning"},
        {"name": "X-Frame-Options", "value": sh.get("x_frame_options", "No"), "status": "Good" if sh.get("x_frame_options") == "Yes" else "Warning"},
        {"name": "X-XSS-Protection", "value": sh.get("x_xss_protection", "No"), "status": "Good" if sh.get("x_xss_protection") == "Yes" else "Warning"},
        {"name": "Broken Links (proxy risk)", "value": live.get("broken_links_detected", 0), "status": "Good" if live.get("broken_links_detected", 0) == 0 else "Warning"},
    ]

    sec = audit_data.get("security_kpis") or {}
    for key in ["critical_vulns", "high_vulns", "medium_vulns", "low_vulns"]:
        if key in sec:
            v = float(sec[key])
            status = "Warning" if v > 0 else "Good"
            rows.append({"name": key.replace("_", " ").title(), "value": v, "status": status})

    return rows[:20]


def _build_accessibility_kpis_table(audit_data: Dict[str, Any], live: Dict[str, Any]) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    rows += [
        {"name": "Alt Text Coverage (%)", "value": live.get("img_alt_coverage_pct", 0), "status": "Good" if live.get("img_alt_coverage_pct", 0) >= 80 else "Warning"},
        {"name": "ARIA Missing (count)", "value": live.get("aria_missing_count", 0), "status": "Good" if live.get("aria_missing_count", 0) <= 10 else "Warning"},
        {"name": "Viewport (mobile)", "value": live.get("viewport_present", "No"), "status": "Good" if live.get("viewport_present") == "Yes" else "Warning"},
        {"name": "Headings Structure (H1)", "value": live.get("h1_count", 0), "status": "Good" if live.get("h1_count", 0) == 1 else "Warning"},
    ]

    acc = audit_data.get("accessibility_kpis") or {}
    for k, v in acc.items():
        label = k.replace("_", " ").title()
        if all(label != r["name"] for r in rows):
            rows.append({"name": label, "value": v, "status": "Good"})

    return rows[:20]


def _build_ux_kpis_table(audit_data: Dict[str, Any], live: Dict[str, Any]) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    rows += [
        {"name": "Mobile-Friendliness (Viewport)", "value": live.get("viewport_present", "No"), "status": "Good" if live.get("viewport_present") == "Yes" else "Warning"},
        {"name": "Broken Links (count)", "value": live.get("broken_links_detected", 0), "status": "Good" if live.get("broken_links_detected", 0) == 0 else "Warning"},
        {"name": "Page Weight (MB)", "value": live.get("page_weight_mb", 0), "status": "Good" if live.get("page_weight_mb", 0) <= 3 else "Warning"},
    ]

    ux = audit_data.get("ux_kpis") or {}
    for k, v in ux.items():
        label = k.replace("_", " ").title()
        if all(label != r["name"] for r in rows):
            rows.append({"name": label, "value": v, "status": "Good"})

    return rows[:20]


def _consolidate_kpis(*sections: List[List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    for sec in sections:
        for row in (sec or []):
            out.append({
                "name": row.get("name", ""),
                "value": row.get("value", ""),
                "status": row.get("status", "")
            })
    return out[:150]


def _collect_ai_recommendations(live: Dict[str, Any], scores: Dict[str, float]) -> List[str]:
    recs: List[str] = []
    def add(msg):
        if msg not in recs:
            recs.append(msg)

    # SEO
    if live.get("title", "") == "" or live.get("meta_description_present") != "Yes":
        add("Add or improve <title> and meta description for better relevance and CTR.")
    if live.get("h1_count", 0) != 1:
        add("Ensure a single, descriptive H1 per page; organize H2/H3 for structure.")
    if live.get("broken_links_detected", 0) > 0:
        add("Fix broken internal links to improve crawlability and UX.")
    if live.get("schema_org_present") != "Yes":
        add("Implement structured data (Schema.org) on key templates.")

    # Performance
    if (isinstance(live.get("ttfb_ms"), (int, float)) and live.get("ttfb_ms") > 800) or (isinstance(live.get("page_weight_mb"), (int, float)) and live.get("page_weight_mb") > 3):
        add("Optimize server TTFB and reduce page weight (image compression, lazy-loading, code splitting).")
    if live.get("compressed") != "Yes":
        add("Enable brotli/gzip compression and HTTP/2/3 where possible.")
    if live.get("cdn", "Unknown") == "Unknown/No signal":
        add("Use a CDN for global delivery and better caching performance.")

    # Security
    sh = live.get("security_headers", {}) or {}
    if sh.get("https_tls") != "Yes" or sh.get("hsts") != "Yes":
        add("Enforce HTTPS with HSTS to prevent protocol downgrade/SSL stripping.")
    if sh.get("csp") != "Yes":
        add("Set a strict Content-Security-Policy (script-src/style-src) to mitigate XSS.")
    if sh.get("x_frame_options") != "Yes":
        add("Set X-Frame-Options or frame-ancestors in CSP to prevent clickjacking.")

    # Accessibility
    if live.get("img_alt_coverage_pct", 100) < 80 or live.get("aria_missing_count", 0) > 10:
        add("Increase alt text coverage and add appropriate ARIA attributes; test with screen readers.")

    # UX
    if live.get("viewport_present") != "Yes":
        add("Add responsive viewport meta for mobile users.")

    return recs[:15]
