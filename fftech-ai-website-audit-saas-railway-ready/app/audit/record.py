# app/audit/record.py

from pathlib import Path
import matplotlib
matplotlib.use('Agg')  # Essential for headless servers like Railway
import matplotlib.pyplot as plt
import pandas as pd
import os
from pptx import Presentation
from pptx.util import Inches

# Ensure the directory exists using absolute path logic
BASE_DIR = Path(__file__).resolve().parent.parent
ASSETS_DIR = BASE_DIR / 'static' / 'generated_assets'
os.makedirs(ASSETS_DIR, exist_ok=True)

def generate_charts(audit_result: dict) -> str:
    """
    Generates a primary breakdown chart for a single audit.
    """
    fig, ax = plt.subplots(figsize=(6, 3))
    breakdown = audit_result.get('breakdown', {})

    cats = list(breakdown.keys())
    vals = [breakdown[k] for k in cats]

    # Matching the blue/green/orange palette used in your frontend CSS
    colors = ['#2E86DE', '#58D68D', '#F5B041', '#EB4D4B', '#A569BD']
    ax.bar(cats, vals, color=colors[:len(cats)])

    ax.set_title('Category Scores Breakdown')
    ax.set_ylim(0, 100)
    ax.set_ylabel('Score')

    chart_path = ASSETS_DIR / 'category_scores.png'
    fig.tight_layout()
    fig.savefig(chart_path, dpi=150)
    plt.close(fig)
    return str(chart_path)

def generate_bar(labels, values, title, filename):
    """
    Generic bar chart generator for competitors or other metrics.
    """
    fig, ax = plt.subplots(figsize=(6, 3))

    # Highlight the base site in blue, competitors in green
    bar_colors = ['#2E86DE'] + ['#58D68D'] * (len(labels) - 1)
    ax.bar(labels, values, color=bar_colors)

    ax.set_title(title)

    numeric_vals = [v for v in values if isinstance(v, (int, float))]
    upper = max(numeric_vals + [100])
    ax.set_ylim(0, upper * 1.1)  # Add 10% padding

    fig.tight_layout()
    out_path = ASSETS_DIR / filename
    fig.savefig(out_path, dpi=150)
    plt.close(fig)
    return str(out_path)

def export_xlsx(audit_result: dict) -> str:
    """
    Exports raw audit data to Excel.
    """
    # Flatten nested dicts like 'breakdown' and 'performance'
    df = pd.json_normalize(audit_result)
    path = ASSETS_DIR / 'audit_summary.xlsx'
    with pd.ExcelWriter(path, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='Audit Results')
    return str(path)

def export_pptx(audit_result: dict, chart_path: str) -> str:
    """
    Generates a professional PowerPoint presentation of the audit.
    """
    prs = Presentation()

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = 'FF Tech - Website Audit Report'
    slide1.placeholders[1].text = (
        f"Overall Score: {audit_result.get('overall_score', 0)}\n"
        f"Grade: {audit_result.get('grade', 'N/A')}"
    )

    # Slide 2: Chart Visualization
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    if os.path.exists(chart_path):
        slide2.shapes.add_picture(chart_path, Inches(1.5), Inches(1.5), height=Inches(4.5))

    path = ASSETS_DIR / 'audit_deck.pptx'
    prs.save(path)
    return str(path)
